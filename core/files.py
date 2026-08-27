# -*- coding: utf-8 -*-
"""
core.files - file extraction, token estimation, context checking.
No streamlit imports.
"""
from core.paths import TEXT_FILE_EXTENSIONS, SUPPORTED_UPLOAD_TYPES
from core.fs import decode_bytes

# Percentage of non-ASCII bytes above which we consider text "mostly non-ASCII"
# (affects the fallback estimate when tiktoken is not available).
_NONASCII_THRESHOLD = 0.2

# Maximum number of tokens allowed for an uploaded file's extracted text.
# Files whose extracted text exceeds this limit are rejected as too large.
MAX_UPLOAD_TOKENS = 500_000


# Maximum extracted-text length for inline attachments. Files larger than this
# are saved into the dialog's files folder; only metadata + a preview are
# passed to the orchestrator (see build_attachment_metadata / build_attachments_context).
MAX_INLINE_UPLOAD_CHARS = 60_000


def get_file_uploader_types() -> list:
    """Return the supported file upload type list."""
    return list(SUPPORTED_UPLOAD_TYPES)


def estimate_tokens(text: str) -> int:
    """Return an accurate token count.

    Uses tiktoken (cl100k_base) if available, otherwise a character-based
    heuristic that accounts for non-ASCII characters (approx 2 tok/char for
    Cyrillic-heavy text vs 0.25 tok/char for pure ASCII).
    Returns at least 1.
    """
    # Try tiktoken first - it is the most accurate for OpenAI-compatible models.
    try:
        import tiktoken
        enc = tiktoken.encoding_for_model("gpt-4")  # cl100k_base
        return max(1, len(enc.encode(text)))
    except ImportError:
        pass

    # Fallback heuristic
    if not text:
        return 1
    ascii_count = sum(1 for ch in text if ord(ch) < 128)
    nonascii_count = len(text) - ascii_count
    # Roughly: 1 token per 4 ASCII chars, 1 token per 2 non-ASCII chars
    # (empirically Cyrillic averages ~2 bytes/char in UTF-8 but tokenizes to ~1.5 tok/char)
    return max(1, ascii_count // 4 + nonascii_count // 2)


def check_upload_tokens(text: str, max_tokens: int = MAX_UPLOAD_TOKENS):
    """Return ``(ok, tokens)`` where ``ok`` is False when the extracted text
    exceeds *max_tokens* tokens.

    Parameters
    ----------
    text : str
        Extracted file content.
    max_tokens : int
        Maximum allowed token count (default: ``MAX_UPLOAD_TOKENS``).

    Returns
    -------
    (bool, int)
        A tuple: ``(tokens <= max_tokens, token_count)``.
    """
    tokens = estimate_tokens(text)
    return tokens <= max_tokens, tokens


def should_store_uploaded_file(text: str, max_chars: int = MAX_INLINE_UPLOAD_CHARS) -> bool:
    """Return True when an uploaded file's extracted text is too large for context."""
    return len(text or "") > max_chars


def build_attachment_metadata(file_name: str, content: str, preview_chars: int = 2000) -> dict:
    """Return an attachment record for one uploaded file.

    Small files are marked ``stored=False`` and keep the full text in
    ``content`` (they will be passed inline). Large files are marked
    ``stored=True``; the full text is still retained in ``content`` so it can
    be written to the dialog's file folder at send time, but only metadata +
    preview are exposed to the orchestrator via ``build_attachments_context``.
    """
    text = content or ""
    return {
        "name": file_name,
        "content": text,
        "path": "",
        "stored": should_store_uploaded_file(text),
        "chars": len(text),
        "tokens": estimate_tokens(text),
        "preview": text[:preview_chars],
    }


def build_attachments_context(attachments: list) -> tuple:
    """Return ``(context_text, joined_names)`` for a list of attachment records.

    Inline attachments (``stored=False``) are embedded in full. Stored
    attachments (``stored=True``) contribute only metadata (saved path, size,
    tokens) plus a small preview - the orchestrator decides when to read the
    full file from disk.
    """
    parts, names = [], []
    for f in attachments or []:
        name = f.get("name", "")
        names.append(name)
        if f.get("stored"):
            preview = f.get("preview", "")
            parts.append(
                f"### {name}\n"
                f"- saved to dialog file folder: {f.get('path', '')}\n"
                f"- size: {f.get('chars', 0)} chars, {f.get('tokens', 0)} tokens\n"
                f"- preview (first {len(preview)} chars):\n"
                f"```\n{preview}\n```"
            )
        else:
            parts.append(f"### {name}\n```\n{f.get('content', '')}\n```")
    return "\n\n".join(parts), ", ".join(names)


def build_saved_files_registry(attachments: list) -> str:
    """Return a compact registry of files saved in the dialog thread folder.

    The registry is small enough to be injected into EVERY user message so
    the orchestrator always knows which files belong to the dialog - even
    when economy mode has truncated the start of the conversation.

    Each entry: ``- name (path, N chars, ~N tokens)``
    """
    if not attachments:
        return ""
    lines = ["Сохранённые файлы диалога:"]
    for f in attachments:
        name = f.get("name", "?")
        path = f.get("path", "") or "?"
        chars = f.get("chars", 0)
        tokens = f.get("tokens", 0)
        lines.append(f"- {name} ({path}, {chars} chars, ~{tokens} tokens)")
    return "\n".join(lines)



def get_model_context_window(skill: dict, services: dict) -> int:
    """Return the context window size for the model referenced by *skill*."""
    svc_name = skill.get("service", "")
    svc      = services.get(svc_name, {})
    model_id = skill.get("model", "")
    for m in svc.get("models", []):
        if isinstance(m, dict) and m.get("id") == model_id:
            return int(m.get("context_window", svc.get("context_window_default", 128000)))
    return int(svc.get("context_window_default", 128000))


def check_context(sys_text: str, user_text: str, file_text: str,
                  skill: dict, services: dict) -> dict:
    """Return a dict describing token usage vs the model's context window."""
    sys_tok  = estimate_tokens(sys_text)
    usr_tok  = estimate_tokens(user_text)
    file_tok = estimate_tokens(file_text)
    total    = sys_tok + usr_tok + file_tok
    limit    = get_model_context_window(skill, services)
    effective_limit = int(limit * 0.9)
    ok              = total <= effective_limit
    excess_tokens   = max(0, total - effective_limit)
    excess_chars    = excess_tokens * 4
    return {
        "sys_tokens":   sys_tok,
        "usr_tokens":   usr_tok,
        "file_tokens":  file_tok,
        "total_tokens": total,
        "limit":        limit,
        "ok":           ok,
        "excess_chars": excess_chars,
    }


def ensure_optional_dependencies() -> list:
    """Return list of optional package names that are not installed."""
    required = {
        "docx":       "python-docx",
        "pptx":       "python-pptx",
        "openpyxl":   "openpyxl",
        "pypdf":      "pypdf",
        "pdfplumber": "pdfplumber",
        "markdown":   "markdown",
        "reportlab":  "reportlab",
    }
    missing = []
    for module_name, package_name in required.items():
        try:
            __import__(module_name)
        except Exception:
            missing.append(package_name)
    return missing


def extract_file_content(uploaded_file) -> str:
    """Extract text content from an uploaded file object (duck-typed: .name, .read())."""
    name = uploaded_file.name.lower()
    raw  = uploaded_file.read()

    if name.endswith(TEXT_FILE_EXTENSIONS):
        return decode_bytes(raw)

    if name.endswith(".docx"):
        try:
            from docx import Document
            import io
            doc   = Document(io.BytesIO(raw))
            lines = [p.text for p in doc.paragraphs if p.text.strip()]
            for tbl in doc.tables:
                if tbl.rows:
                    hrow = tbl.rows[0]
                    lines += [
                        "",
                        "| " + " | ".join(c.text for c in hrow.cells) + " |",
                        "| " + " | ".join("---" for _ in hrow.cells)  + " |",
                    ]
                    for row in tbl.rows[1:]:
                        lines.append("| " + " | ".join(c.text for c in row.cells) + " |")
            return "\n".join(lines)
        except Exception as e:
            raise RuntimeError(f"docx: {e}")

    if name.endswith(".pptx"):
        try:
            from pptx import Presentation
            import io
            prs   = Presentation(io.BytesIO(raw))
            lines = []
            for idx, slide in enumerate(prs.slides, 1):
                lines.append(f"\n## Slide {idx}")
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                lines.append(para.text.strip())
                    if shape.has_table:
                        tbl  = shape.table
                        cols = len(tbl.columns)
                        hdr  = [tbl.cell(0, c).text for c in range(cols)]
                        lines += [
                            "| " + " | ".join(hdr) + " |",
                            "| " + " | ".join("---" for _ in hdr) + " |",
                        ]
                        for r in range(1, len(tbl.rows)):
                            lines.append("| " + " | ".join(tbl.cell(r, c).text for c in range(cols)) + " |")
            return "\n".join(lines)
        except Exception as e:
            raise RuntimeError(f"pptx: {e}")

    if name.endswith(".xlsx"):
        try:
            import openpyxl, io
            wb    = openpyxl.load_workbook(io.BytesIO(raw), data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"\n## Sheet: {sheet.title}")
                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    continue
                header = [str(v) if v is not None else "" for v in rows[0]]
                lines += [
                    "| " + " | ".join(header) + " |",
                    "| " + " | ".join("---" for _ in header) + " |",
                ]
                for row in rows[1:]:
                    lines.append("| " + " | ".join(str(v) if v is not None else "" for v in row) + " |")
            return "\n".join(lines)
        except Exception as e:
            raise RuntimeError(f"xlsx: {e}")

    if name.endswith(".pdf"):
        try:
            import io
            try:
                import pdfplumber
                lines = []
                with pdfplumber.open(io.BytesIO(raw)) as pdf:
                    for i, page in enumerate(pdf.pages, 1):
                        text = page.extract_text() or ""
                        if text.strip():
                            lines.append(f"### Page {i}\n{text.strip()}")
                        tables = page.extract_tables()
                        for tbl in tables:
                            if not tbl:
                                continue
                            hdr = [str(c) if c else "" for c in tbl[0]]
                            lines.append("| " + " | ".join(hdr) + " |")
                            lines.append("| " + " | ".join("---" for _ in hdr) + " |")
                            for row in tbl[1:]:
                                lines.append("| " + " | ".join(str(c) if c else "" for c in row) + " |")
                return "\n\n".join(lines)
            except ImportError:
                pass
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(raw))
            lines = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    lines.append(f"### Page {i}\n{text.strip()}")
            return "\n\n".join(lines)
        except Exception as e:
            raise RuntimeError(f"pdf: {e}")

    raise RuntimeError("Unsupported file format")
# SPDX-FileCopyrightText: 2026 SagaAI Platform, Deinekin T.V.
# SPDX-License-Identifier: MIT
